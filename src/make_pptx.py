# src/make_pptx.py
import os
from pptx import Presentation
from pptx.util import Inches

# ----------------------------------------------------------------
# Setup paths
# ----------------------------------------------------------------
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR = os.path.join(BASE_DIR, "outputs", "figures")
OUT_DIR = os.path.join(BASE_DIR, "outputs")

os.makedirs(OUT_DIR, exist_ok=True)

# ----------------------------------------------------------------
# Create presentation
# ----------------------------------------------------------------
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "EDA: Local Weather-Based Disease Prediction"
slide.placeholders[1].text = "Dataset: Weather-related disease prediction.csv\nPrepared by: Aniket Bansod"

# Slide 1: Data Snapshot
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Data snapshot & missing values"
slide.placeholders[1].text = "Explain file shape, missing values and data types.\nSee table: outputs/tables/missing.csv"

# Slide 2: Target distribution
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Target distribution"
img_path = os.path.join(FIG_DIR, "target_distribution.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Slide 3: Symptom frequencies
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Top symptom frequencies"
img_path = os.path.join(FIG_DIR, "symptom_freq_top30.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.2), width=Inches(8))

# Slide 4: Symptom co-occurrence
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Symptom co-occurrence"
img_path = os.path.join(FIG_DIR, "symptom_clustermap.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(8))

# Slide 5: Weather vs disease
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Weather distribution by disease"
img_path = os.path.join(FIG_DIR, "temp_by_disease_top6.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.2), width=Inches(8))

# Slide 6: Statistical tests & summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Statistical tests & key takeaways"
slide.placeholders[1].text = (
    "1) Chi-square: symptom ~ disease (p < 0.05)\n"
    "2) ANOVA/Kruskal: weather ~ disease (p < 0.05)\n"
    "3) Next steps: preprocessing & predictive modeling"
)

# ----------------------------------------------------------------
# Save
# ----------------------------------------------------------------
out_file = os.path.join(OUT_DIR, "EDA_presentation.pptx")
prs.save(out_file)
print(f"✅ Saved presentation at: {out_file}")
